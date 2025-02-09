from __future__ import annotations

import re
from contextlib import suppress
from html import escape
from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import TYPE_CHECKING

import html_to_markdown
import pptx
import pypdfium2
from anyio import Path as AsyncPath
from pptx.enum.shapes import MSO_SHAPE_TYPE
from xlsx2csv import Xlsx2csv

from kreuzberg._pandoc import process_content, process_file
from kreuzberg._string import normalize_spaces, safe_decode
from kreuzberg._sync import run_sync
from kreuzberg._tesseract import batch_process_images
from kreuzberg.exceptions import ParsingError

if TYPE_CHECKING:  # pragma: no cover
    from PIL.Image import Image


async def convert_pdf_to_images(file_path: Path) -> list[Image]:
    """Convert a PDF file to images.

    Args:
        file_path: The path to the PDF file.

    Raises:
        ParsingError: If the PDF file could not be converted to images.

    Returns:
        A list of Pillow Images.
    """
    try:
        resolved_path = str(await AsyncPath(file_path).resolve())
        pdf = await run_sync(pypdfium2.PdfDocument, resolved_path)
        return [page.render(scale=2.0).to_pil() for page in pdf]
    except pypdfium2.PdfiumError as e:
        raise ParsingError(
            "Could not convert PDF to images", context={"file_path": str(file_path), "error": str(e)}
        ) from e


async def extract_pdf_with_tesseract(file_path: Path) -> str:
    """Extract text from a scanned PDF file using pytesseract.

    Args:
        file_path: The path to the PDF file.

    Returns:
        The extracted text.
    """
    images = await convert_pdf_to_images(file_path)
    ocr_results = await batch_process_images(images)
    return normalize_spaces("\n".join(ocr_results))


async def extract_pdf_with_pdfium2(file_path: Path) -> str:
    """Extract text from a searchable PDF file using pypdfium2.

    Args:
        file_path: The path to the PDF file.

    Raises:
        ParsingError: If the text could not be extracted from the PDF file.

    Returns:
        The extracted text.
    """
    try:
        resolved_path = str(await AsyncPath(file_path).resolve())
        document = await run_sync(pypdfium2.PdfDocument, resolved_path)
        text = "\n".join(page.get_textpage().get_text_bounded() for page in document)
        return normalize_spaces(text)
    except pypdfium2.PdfiumError as e:
        raise ParsingError(
            "Could not extract text from PDF file", context={"file_path": str(file_path), "error": str(e)}
        ) from e


async def extract_pdf_file(file_path: Path, force_ocr: bool = False) -> str:
    """Extract text from a PDF file.

    Args:
        file_path: The path to the PDF file.
        force_ocr: Whether or not to force OCR on PDF files that have a text layer. Default = false.

    Returns:
        The extracted text.
    """
    if not force_ocr and (content := await extract_pdf_with_pdfium2(file_path)):
        return normalize_spaces(content)

    return await extract_pdf_with_tesseract(file_path)


async def extract_content_with_pandoc(file_data: bytes, mime_type: str) -> str:
    """Extract text using pandoc.

    Args:
        file_data: The content of the file.
        mime_type: The mime type of the file.

    Returns:
        The extracted text.
    """
    result = await process_content(file_data, mime_type=mime_type)
    return normalize_spaces(result.content)


async def extract_file_with_pandoc(file_path: Path | str, mime_type: str) -> str:
    """Extract text using pandoc.

    Args:
        file_path: The path to the file.
        mime_type: The mime type of the file.

    Returns:
        The extracted text.
    """
    resolved_path = str(await AsyncPath(file_path).resolve())
    result = await process_file(resolved_path, mime_type=mime_type)
    return normalize_spaces(result.content)


async def extract_pptx_file(file_path_or_contents: Path | bytes) -> str:
    """Extract text from a PPTX file.

    Notes:
        This function is based on code vendored from `markitdown`, which has an MIT license as well.

    Args:
        file_path_or_contents: The path to the PPTX file or its contents as bytes.

    Returns:
        The extracted text content
    """
    md_content = ""
    file_contents = (
        file_path_or_contents
        if isinstance(file_path_or_contents, bytes)
        else await AsyncPath(file_path_or_contents).read_bytes()
    )
    presentation = pptx.Presentation(BytesIO(file_contents))

    for index, slide in enumerate(presentation.slides):
        md_content += f"\n\n<!-- Slide number: {index + 1} -->\n"

        title = slide.shapes.title

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
                shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image")
            ):
                alt_text = ""
                with suppress(AttributeError):
                    # access non-visual properties
                    alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")  # noqa: SLF001

                filename = re.sub(r"\W", "", shape.name) + ".jpg"
                md_content += f"\n![{alt_text if alt_text else shape.name}]({filename})\n"

            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                html_table = "<table>"
                first_row = True

                for row in shape.table.rows:
                    html_table += "<tr>"

                    for cell in row.cells:
                        tag = "th" if first_row else "td"
                        html_table += f"<{tag}>{escape(cell.text)}</{tag}>"

                    html_table += "</tr>"
                    first_row = False

                html_table += "</table>"
                md_content += "\n" + html_table + "\n"

            elif shape.has_text_frame:
                md_content += "# " + shape.text.lstrip() + "\n" if shape == title else shape.text + "\n"

        md_content = md_content.strip()
        if slide.has_notes_slide:
            md_content += "\n\n### Notes:\n"
            notes_frame = slide.notes_slide.notes_text_frame

            if notes_frame is not None:
                md_content += notes_frame.text

            md_content = md_content.strip()

    return normalize_spaces(md_content)


async def extract_xlsx_file(file_path_or_contents: Path | bytes) -> str:
    """Extract text from an XLSX file by converting it to CSV and then to markdown.

    Args:
        file_path_or_contents: The path to the XLSX file or its contents as bytes.

    Returns:
        The extracted text content.

    Raises:
        ParsingError: If the XLSX file could not be parsed.
    """
    try:
        with NamedTemporaryFile(suffix=".xlsx") as xlsx_file, NamedTemporaryFile(suffix=".csv") as csv_file:
            if isinstance(file_path_or_contents, bytes):
                xlsx_file.write(file_path_or_contents)
                xlsx_file.flush()
                xlsx_path = xlsx_file.name
            else:
                xlsx_path = str(await AsyncPath(file_path_or_contents).resolve())

            await run_sync(Xlsx2csv(xlsx_path).convert, csv_file.name)
            result = await process_file(csv_file.name, mime_type="text/csv")
            return normalize_spaces(result.content)
    except Exception as e:
        raise ParsingError(
            "Could not extract text from XLSX file",
            context={
                "error": str(e),
                "file_path": str(file_path_or_contents) if isinstance(file_path_or_contents, Path) else None,
            },
        ) from e


async def extract_html_string(file_path_or_contents: Path | bytes) -> str:
    """Extract text from an HTML string.

    Args:
        file_path_or_contents: The HTML content.

    Returns:
        The extracted text content.
    """
    content = (
        safe_decode(file_path_or_contents)
        if isinstance(file_path_or_contents, bytes)
        else await AsyncPath(file_path_or_contents).read_text()
    )
    return normalize_spaces(await run_sync(html_to_markdown.convert_to_markdown, content))
