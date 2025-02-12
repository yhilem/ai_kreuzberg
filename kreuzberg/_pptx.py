from __future__ import annotations

import re
from contextlib import suppress
from html import escape
from io import BytesIO
from typing import TYPE_CHECKING

import pptx
from anyio import Path as AsyncPath
from pptx.enum.shapes import MSO_SHAPE_TYPE

from kreuzberg import ExtractionResult
from kreuzberg._mime_types import MARKDOWN_MIME_TYPE
from kreuzberg._string import normalize_spaces

if TYPE_CHECKING:
    from pathlib import Path


async def extract_pptx_file_content(file_path_or_contents: Path | bytes) -> ExtractionResult:
    """Extract text from a PPTX file.

    Notes:
        This function is based on code vendored from `markitdown`, which has an MIT license as well.

    Args:
        file_path_or_contents: The path to the PPTX file or its contents as bytes.

    Returns:
        The extracted text content
    """
    md_content = ""
    file_contents = (
        file_path_or_contents
        if isinstance(file_path_or_contents, bytes)
        else await AsyncPath(file_path_or_contents).read_bytes()
    )
    presentation = pptx.Presentation(BytesIO(file_contents))

    for index, slide in enumerate(presentation.slides):
        md_content += f"\n\n<!-- Slide number: {index + 1} -->\n"

        title = slide.shapes.title

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
                shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image")
            ):
                alt_text = ""
                with suppress(AttributeError):
                    # access non-visual properties
                    alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")  # noqa: SLF001

                filename = re.sub(r"\W", "", shape.name) + ".jpg"
                md_content += f"\n![{alt_text if alt_text else shape.name}]({filename})\n"

            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                html_table = "<table>"
                first_row = True

                for row in shape.table.rows:
                    html_table += "<tr>"

                    for cell in row.cells:
                        tag = "th" if first_row else "td"
                        html_table += f"<{tag}>{escape(cell.text)}</{tag}>"

                    html_table += "</tr>"
                    first_row = False

                html_table += "</table>"
                md_content += "\n" + html_table + "\n"

            elif shape.has_text_frame:
                md_content += "# " + shape.text.lstrip() + "\n" if shape == title else shape.text + "\n"

        md_content = md_content.strip()
        if slide.has_notes_slide:
            md_content += "\n\n### Notes:\n"
            notes_frame = slide.notes_slide.notes_text_frame

            if notes_frame is not None:
                md_content += notes_frame.text

            md_content = md_content.strip()

    return ExtractionResult(content=normalize_spaces(md_content), mime_type=MARKDOWN_MIME_TYPE, metadata={})
